import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pdf2image import convert_from_bytes
import PyPDF2
import tempfile
import os

def convert_pdf_to_ppt(pdf_file, title):
    """
    Convert a PDF file to a PowerPoint presentation.
    
    Args:
        pdf_file: PDF file as a BytesIO object
        title: Title for the presentation
        
    Returns:
        PowerPoint file as bytes
    """
    # Create a presentation
    prs = Presentation()
    
    # Update the properties
    prs.core_properties.title = title
    
    # Read PDF content
    pdf_reader = PyPDF2.PdfReader(pdf_file)
    num_pages = len(pdf_reader.pages)
    pdf_file.seek(0)  # Reset file pointer
    
    # Create a title slide
    title_slide_layout = prs.slide_layouts[0]  # Title slide layout
    title_slide = prs.slides.add_slide(title_slide_layout)
    
    # Set the title
    title_shape = title_slide.shapes.title
    title_shape.text = title
    
    # Subtitle (optional)
    if title_slide.placeholders[1].has_text_frame:
        subtitle = title_slide.placeholders[1]
        subtitle.text = f"Converted PDF Presentation ({num_pages} pages)"
    
    # Convert each PDF page to an image and add it to the presentation
    with tempfile.TemporaryDirectory() as temp_dir:
        # Convert PDF pages to images
        images = convert_from_bytes(pdf_file.read(), dpi=300)
        pdf_file.seek(0)  # Reset file pointer
        
        # Add each page as a slide
        blank_slide_layout = prs.slide_layouts[6]  # Blank slide layout
        
        for i, image in enumerate(images):
            # Save the image temporarily
            img_path = os.path.join(temp_dir, f"page_{i+1}.png")
            image.save(img_path, "PNG")
            
            # Create a new slide
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Add the image to fill most of the slide
            slide.shapes.add_picture(
                img_path, 
                left=Inches(0.5), 
                top=Inches(0.5),
                width=Inches(9),  # Adjust as needed
                height=None      # Maintain aspect ratio
            )
            
            # Add page number
            page_num = slide.shapes.add_textbox(
                left=Inches(9), 
                top=Inches(6.5),
                width=Inches(1),
                height=Inches(0.5)
            )
            tf = page_num.text_frame
            p = tf.add_paragraph()
            p.text = f"Page {i+1}"
            p.font.size = Pt(12)
    
    # Save the presentation to a BytesIO object
    ppt_bytes = io.BytesIO()
    prs.save(ppt_bytes)
    ppt_bytes.seek(0)
    
    return ppt_bytes.getvalue()
